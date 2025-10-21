"""
Сервис для экспорта контента в различные форматы
"""
from typing import Dict, Any
from io import BytesIO
import logging

from backend.models.domain import Course, Module

logger = logging.getLogger(__name__)


class ExportService:
    """Сервис для экспорта курсов, модулей и уроков в различные форматы"""
    
    # ========== ЭКСПОРТ КУРСА ==========
    
    @staticmethod
    def export_course_markdown(course: Course) -> str:
        """Генерирует Markdown для всего курса"""
        md = f"# {course.course_title}\n\n"
        md += f"**Целевая аудитория:** {course.target_audience}\n\n"
        
        if course.duration_weeks:
            md += f"**Длительность:** {course.duration_weeks} недель\n\n"
        if course.duration_hours:
            md += f"**Всего часов:** {course.duration_hours}\n\n"
        
        md += "---\n\n"
        
        for module in course.modules:
            md += f"## Модуль {module.module_number}: {module.module_title}\n\n"
            md += f"**Цель модуля:** {module.module_goal}\n\n"
            
            for i, lesson in enumerate(module.lessons, 1):
                md += f"### {i}. {lesson.lesson_title}\n\n"
                md += f"**Цель урока:** {lesson.lesson_goal}\n\n"
                md += f"**Формат:** {lesson.format} | **Время:** {lesson.estimated_time_minutes} мин\n\n"
                md += f"**План контента:**\n"
                for item in lesson.content_outline:
                    md += f"- {item}\n"
                md += f"\n**Оценка:** {lesson.assessment}\n\n"
            
            md += "---\n\n"
        
        return md
    
    @staticmethod
    def export_course_text(course: Course) -> str:
        """Генерирует текстовый формат курса"""
        txt = f"{course.course_title}\n"
        txt += "=" * len(course.course_title) + "\n\n"
        txt += f"Целевая аудитория: {course.target_audience}\n"
        
        if course.duration_weeks:
            txt += f"Длительность: {course.duration_weeks} недель\n"
        if course.duration_hours:
            txt += f"Всего часов: {course.duration_hours}\n"
        
        txt += "\n" + "-" * 80 + "\n\n"
        
        for module in course.modules:
            txt += f"\nМОДУЛЬ {module.module_number}: {module.module_title}\n"
            txt += "-" * 40 + "\n"
            txt += f"Цель модуля: {module.module_goal}\n\n"
            
            for i, lesson in enumerate(module.lessons, 1):
                txt += f"  {i}. {lesson.lesson_title}\n"
                txt += f"     Цель: {lesson.lesson_goal}\n"
                txt += f"     Формат: {lesson.format} | Время: {lesson.estimated_time_minutes} мин\n"
                txt += f"     План контента:\n"
                for item in lesson.content_outline:
                    txt += f"       - {item}\n"
                txt += f"     Оценка: {lesson.assessment}\n\n"
            
            txt += "\n"
        
        return txt
    
    @staticmethod
    def export_course_html(course: Course) -> str:
        """Генерирует HTML для всего курса"""
        html = f"""<!DOCTYPE html>
<html lang="ru">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{course.course_title}</title>
    <style>
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            max-width: 1200px;
            margin: 0 auto;
            padding: 20px;
            background-color: #f5f5f5;
            line-height: 1.6;
        }}
        .header {{
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 40px;
            border-radius: 10px;
            margin-bottom: 30px;
        }}
        .header h1 {{
            margin: 0 0 10px 0;
        }}
        .meta {{
            display: flex;
            gap: 20px;
            flex-wrap: wrap;
        }}
        .meta-item {{
            background: rgba(255,255,255,0.2);
            padding: 5px 15px;
            border-radius: 5px;
        }}
        .module {{
            background: white;
            padding: 30px;
            margin-bottom: 20px;
            border-radius: 10px;
            box-shadow: 0 2px 5px rgba(0,0,0,0.1);
        }}
        .module h2 {{
            color: #667eea;
            border-bottom: 2px solid #667eea;
            padding-bottom: 10px;
            margin-top: 0;
        }}
        .module-goal {{
            background: #f0f4ff;
            padding: 15px;
            border-left: 4px solid #667eea;
            margin: 15px 0;
        }}
        .lesson {{
            margin: 20px 0;
            padding: 20px;
            background: #fafafa;
            border-radius: 5px;
        }}
        .lesson h3 {{
            color: #764ba2;
            margin-top: 0;
        }}
        .lesson-meta {{
            color: #666;
            margin: 10px 0;
        }}
        .lesson-meta span {{
            background: #e0e0e0;
            padding: 3px 10px;
            border-radius: 3px;
            margin-right: 10px;
        }}
        .content-outline {{
            background: white;
            padding: 15px;
            border-radius: 5px;
            margin: 10px 0;
        }}
        .content-outline ul {{
            margin: 5px 0;
            padding-left: 20px;
        }}
        .assessment {{
            background: #fff9e6;
            padding: 10px;
            border-left: 3px solid #ffd700;
            margin-top: 10px;
        }}
        @media print {{
            body {{ background: white; }}
            .module {{ page-break-inside: avoid; }}
        }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{course.course_title}</h1>
        <div class="meta">
            <div class="meta-item">👥 {course.target_audience}</div>
"""
        
        if course.duration_weeks:
            html += f'            <div class="meta-item">📅 {course.duration_weeks} недель</div>\n'
        if course.duration_hours:
            html += f'            <div class="meta-item">⏱️ {course.duration_hours} часов</div>\n'
        
        html += """        </div>
    </div>
"""
        
        for module in course.modules:
            html += f"""
    <div class="module">
        <h2>Модуль {module.module_number}: {module.module_title}</h2>
        <div class="module-goal">
            <strong>🎯 Цель модуля:</strong> {module.module_goal}
        </div>
"""
            
            for i, lesson in enumerate(module.lessons, 1):
                html += f"""
        <div class="lesson">
            <h3>{i}. {lesson.lesson_title}</h3>
            <p><strong>Цель урока:</strong> {lesson.lesson_goal}</p>
            <div class="lesson-meta">
                <span>📚 {lesson.format}</span>
                <span>⏱️ {lesson.estimated_time_minutes} мин</span>
            </div>
            <div class="content-outline">
                <strong>План контента:</strong>
                <ul>
"""
                for item in lesson.content_outline:
                    html += f"                    <li>{item}</li>\n"
                
                html += f"""                </ul>
            </div>
            <div class="assessment">
                <strong>✅ Оценка:</strong> {lesson.assessment}
            </div>
        </div>
"""
            
            html += "    </div>\n"
        
        html += """
</body>
</html>"""
        
        return html
    
    @staticmethod
    def export_course_pptx(course: Course) -> bytes:
        """Генерирует PowerPoint презентацию курса"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Титульный слайд
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(102, 126, 234)
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.text = course.course_title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.alignment = PP_ALIGN.CENTER
            
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_text = f"Целевая аудитория: {course.target_audience}"
            if course.duration_weeks:
                subtitle_text += f" | {course.duration_weeks} недель"
            subtitle_frame.text = subtitle_text
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.alignment = PP_ALIGN.CENTER
            
            # Слайд структуры курса
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Структура курса"
            
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = f"Курс состоит из {len(course.modules)} модулей:"
            
            for module in course.modules:
                p = tf.add_paragraph()
                p.text = f"Модуль {module.module_number}: {module.module_title}"
                p.level = 1
                p.font.size = Pt(18)
            
            # Слайды для каждого модуля
            for module in course.modules:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = f"Модуль {module.module_number}: {module.module_title}"
                
                content = slide.placeholders[1]
                tf = content.text_frame
                tf.text = f"Цель: {module.module_goal}"
                
                p = tf.add_paragraph()
                p.text = f"\nУроков в модуле: {len(module.lessons)}"
                p.font.size = Pt(18)
            
            pptx_bytes_io = BytesIO()
            prs.save(pptx_bytes_io)
            pptx_bytes_io.seek(0)
            
            return pptx_bytes_io.getvalue()
            
        except Exception as e:
            logger.error(f"Ошибка генерации PPTX курса: {e}")
            raise
    
    # ========== ЭКСПОРТ МОДУЛЯ (ДЕТАЛЬНЫЙ КОНТЕНТ) ==========
    
    @staticmethod
    def export_module_markdown(course: Course, module: Module, content_data: dict) -> str:
        """Генерирует Markdown для детального контента модуля"""
        md = f"# {course.course_title}\n\n"
        md += f"## Модуль {module.module_number}: {module.module_title}\n\n"
        md += f"**Цель модуля:** {module.module_goal}\n\n"
        md += "---\n\n"
        
        lectures = content_data.get('lectures', [])
        for i, lecture in enumerate(lectures, 1):
            md += f"### Лекция {i}: {lecture.get('lecture_title', 'Без названия')}\n\n"
            
            if lecture.get('learning_objectives'):
                md += "**Цели обучения:**\n"
                for obj in lecture['learning_objectives']:
                    md += f"- {obj}\n"
                md += "\n"
            
            if lecture.get('key_takeaways'):
                md += "**Ключевые выводы:**\n"
                for key in lecture['key_takeaways']:
                    md += f"- {key}\n"
                md += "\n"
            
            slides = lecture.get('slides', [])
            for j, slide in enumerate(slides, 1):
                slide_title = slide.get('slide_title') or slide.get('title', 'Без названия')
                slide_content = slide.get('slide_content') or slide.get('content', '')
                
                md += f"#### Слайд {j}: {slide_title}\n\n"
                
                if slide_content:
                    md += f"{slide_content}\n\n"
                
                if slide.get('code_example'):
                    md += f"```python\n{slide['code_example']}\n```\n\n"
                
                if slide.get('visual_description'):
                    md += f"📊 **Визуализация:** {slide['visual_description']}\n\n"
                
                if slide.get('notes'):
                    md += f"📝 _Заметки преподавателя: {slide['notes']}_\n\n"
            
            md += "---\n\n"
        
        return md
    
    @staticmethod  
    def export_module_html(course: Course, module: Module, content_data: dict) -> str:
        """Генерирует HTML для детального контента модуля"""
        html = f"""<!DOCTYPE html>
<html lang="ru">
<head>
    <meta charset="UTF-8">
    <title>{module.module_title} - Детальный контент</title>
    <style>
        body {{ font-family: Arial, sans-serif; max-width: 1000px; margin: 0 auto; padding: 20px; }}
        .header {{ background: #667eea; color: white; padding: 30px; border-radius: 8px; }}
        .lecture {{ background: #f9f9f9; padding: 20px; margin: 20px 0; border-radius: 8px; }}
        .slide {{ background: white; padding: 15px; margin: 10px 0; border-left: 4px solid #667eea; }}
        .code {{ background: #282c34; color: #abb2bf; padding: 15px; border-radius: 4px; overflow-x: auto; }}
        .visual {{ background: #e8f5e9; padding: 10px; border-left: 4px solid #4caf50; margin-top: 10px; }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{course.course_title}</h1>
        <h2>Модуль {module.module_number}: {module.module_title}</h2>
        <p><strong>Цель:</strong> {module.module_goal}</p>
    </div>
"""
        
        lectures = content_data.get('lectures', [])
        for i, lecture in enumerate(lectures, 1):
            html += f"""
    <div class="lecture">
        <h3>Лекция {i}: {lecture.get('lecture_title', 'Без названия')}</h3>
"""
            
            if lecture.get('learning_objectives'):
                html += """        <div style="background: #f0f4ff; padding: 10px; margin: 10px 0; border-left: 3px solid #667eea;">
            <strong>Цели обучения:</strong>
            <ul>
"""
                for obj in lecture['learning_objectives']:
                    html += f"                <li>{obj}</li>\n"
                html += """            </ul>
        </div>
"""
            
            if lecture.get('key_takeaways'):
                html += """        <div style="background: #fff9e6; padding: 10px; margin: 10px 0; border-left: 3px solid #ffd700;">
            <strong>Ключевые выводы:</strong>
            <ul>
"""
                for key in lecture['key_takeaways']:
                    html += f"                <li>{key}</li>\n"
                html += """            </ul>
        </div>
"""
            
            slides = lecture.get('slides', [])
            for j, slide in enumerate(slides, 1):
                slide_title = slide.get('slide_title') or slide.get('title', 'Без названия')
                slide_content = slide.get('slide_content') or slide.get('content', '')
                
                html += f"""
        <div class="slide">
            <h4>Слайд {j}: {slide_title}</h4>
"""
                
                if slide_content:
                    html += f"            <p>{slide_content}</p>\n"
                
                if slide.get('code_example'):
                    html += f"""
            <div class="code">
                <pre><code>{slide['code_example']}</code></pre>
            </div>
"""
                
                if slide.get('visual_description'):
                    html += f"""
            <div class="visual">
                <strong>📊 Визуализация:</strong> {slide['visual_description']}
            </div>
"""
                
                if slide.get('notes'):
                    html += f"""
            <div style="background: #f0f0f0; padding: 8px; margin-top: 10px; font-size: 0.9em; font-style: italic;">
                <strong>📝 Заметки:</strong> {slide['notes']}
            </div>
"""
                
                html += "        </div>\n"
            
            html += "    </div>\n"
        
        html += """
</body>
</html>"""
        
        return html
    
    @staticmethod
    def export_module_pptx(course: Course, module: Module, content_data: dict) -> bytes:
        """Генерирует PowerPoint для детального контента модуля"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Титульный слайд модуля
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(102, 126, 234)
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
            title_frame = title_box.text_frame
            title_frame.text = f"{course.course_title}\n\nМодуль {module.module_number}: {module.module_title}"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.alignment = PP_ALIGN.CENTER
            
            goal_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
            goal_frame = goal_box.text_frame
            goal_frame.text = f"Цель: {module.module_goal}"
            goal_para = goal_frame.paragraphs[0]
            goal_para.font.size = Pt(18)
            goal_para.font.color.rgb = RGBColor(255, 255, 255)
            goal_para.alignment = PP_ALIGN.CENTER
            
            # Генерируем слайды из контента
            lectures = content_data.get('lectures', [])
            
            for lecture_idx, lecture in enumerate(lectures, 1):
                lecture_title = lecture.get('lecture_title', f'Лекция {lecture_idx}')
                
                # Заголовочный слайд лекции
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title = slide.shapes.title
                title.text = f"Лекция {lecture_idx}: {lecture_title}"
                
                # Цели обучения
                if lecture.get('learning_objectives'):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = "Цели обучения"
                    
                    content_shape = slide.placeholders[1]
                    tf = content_shape.text_frame
                    tf.text = ""
                    
                    for obj in lecture['learning_objectives']:
                        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
                        p.text = obj
                        p.level = 0
                        p.font.size = Pt(20)
                        p.space_before = Pt(6)
                
                # Слайды лекции
                slides_data = lecture.get('slides', [])
                for slide_data in slides_data:
                    slide_title = slide_data.get('slide_title') or slide_data.get('title', 'Без названия')
                    slide_content = slide_data.get('slide_content') or slide_data.get('content', '')
                    code_example = slide_data.get('code_example')
                    
                    if code_example:
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        title = slide.shapes.title
                        title.text = slide_title
                        
                        content_shape = slide.placeholders[1]
                        tf = content_shape.text_frame
                        tf.word_wrap = True
                        
                        if slide_content:
                            tf.text = slide_content.replace('\\n', '\n')
                            tf.paragraphs[0].font.size = Pt(16)
                        
                        code_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(3))
                        code_frame = code_box.text_frame
                        code_frame.text = code_example.replace('\\n', '\n')
                        code_para = code_frame.paragraphs[0]
                        code_para.font.name = 'Courier New'
                        code_para.font.size = Pt(14)
                        code_box.fill.solid()
                        code_box.fill.fore_color.rgb = RGBColor(40, 44, 52)
                        code_para.font.color.rgb = RGBColor(171, 178, 191)
                    else:
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        title = slide.shapes.title
                        title.text = slide_title
                        
                        if slide_content:
                            content_shape = slide.placeholders[1]
                            tf = content_shape.text_frame
                            tf.word_wrap = True
                            
                            content_parts = slide_content.replace('\\n', '\n').split('\n')
                            
                            for i, part in enumerate(content_parts):
                                part = part.strip()
                                if not part:
                                    continue
                                
                                if i == 0:
                                    tf.text = part
                                    tf.paragraphs[0].font.size = Pt(18)
                                    tf.paragraphs[0].space_after = Pt(12)
                                else:
                                    p = tf.add_paragraph()
                                    p.text = part
                                    p.font.size = Pt(16)
                                    p.space_before = Pt(6)
                                    p.space_after = Pt(6)
                                    
                                    if part.startswith('- '):
                                        p.text = part[2:]
                                        p.level = 1
                
                # Ключевые выводы лекции
                if lecture.get('key_takeaways'):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = f"Ключевые выводы: {lecture_title}"
                    
                    content_shape = slide.placeholders[1]
                    tf = content_shape.text_frame
                    tf.text = ""
                    
                    for key in lecture['key_takeaways']:
                        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
                        p.text = key
                        p.level = 0
                        p.font.size = Pt(20)
                        p.font.bold = True
                        p.space_before = Pt(12)
            
            pptx_bytes_io = BytesIO()
            prs.save(pptx_bytes_io)
            pptx_bytes_io.seek(0)
            
            return pptx_bytes_io.getvalue()
            
        except Exception as e:
            logger.error(f"Ошибка генерации PPTX модуля: {e}")
            raise
    
    # ========== ЭКСПОРТ УРОКА (ДЕТАЛЬНЫЙ КОНТЕНТ) ==========
    
    @staticmethod
    def export_lesson_markdown(course: Course, module: Module, lesson, content_data: dict) -> str:
        """Генерирует Markdown для детального контента урока"""
        md = f"# {course.course_title}\n\n"
        md += f"## Модуль {module.module_number}: {module.module_title}\n\n"
        md += f"### Урок: {lesson.lesson_title}\n\n"
        md += f"**Цель урока:** {lesson.lesson_goal}\n\n"
        md += f"**Формат:** {lesson.format} | **Время:** {lesson.estimated_time_minutes} мин\n\n"
        md += "---\n\n"
        
        if content_data.get('learning_objectives'):
            md += "**Цели обучения:**\n"
            for obj in content_data['learning_objectives']:
                md += f"- {obj}\n"
            md += "\n"
        
        slides = content_data.get('slides', [])
        for j, slide in enumerate(slides, 1):
            slide_title = slide.get('slide_title') or slide.get('title', 'Без названия')
            slide_content = slide.get('slide_content') or slide.get('content', '')
            
            md += f"#### Слайд {j}: {slide_title}\n\n"
            
            if slide_content:
                md += f"{slide_content}\n\n"
            
            if slide.get('code_example'):
                md += f"```python\n{slide['code_example']}\n```\n\n"
            
            if slide.get('notes'):
                md += f"📝 _{slide['notes']}_\n\n"
        
        if content_data.get('key_takeaways'):
            md += "---\n\n**Ключевые выводы:**\n"
            for key in content_data['key_takeaways']:
                md += f"- {key}\n"
            md += "\n"
        
        return md
    
    @staticmethod
    def export_lesson_html(course: Course, module: Module, lesson, content_data: dict) -> str:
        """Генерирует HTML для детального контента урока"""
        html = f"""<!DOCTYPE html>
<html lang="ru">
<head>
    <meta charset="UTF-8">
    <title>{lesson.lesson_title}</title>
    <style>
        body {{ font-family: Arial, sans-serif; max-width: 900px; margin: 0 auto; padding: 20px; }}
        .header {{ background: #667eea; color: white; padding: 25px; border-radius: 8px; margin-bottom: 20px; }}
        .objectives {{ background: #f0f4ff; padding: 15px; margin: 15px 0; border-left: 4px solid #667eea; }}
        .slide {{ background: #f9f9f9; padding: 20px; margin: 15px 0; border-radius: 8px; }}
        .slide h3 {{ color: #667eea; margin-top: 0; }}
        .code {{ background: #282c34; color: #abb2bf; padding: 15px; border-radius: 4px; overflow-x: auto; margin: 10px 0; }}
        .notes {{ background: #f0f0f0; padding: 10px; margin-top: 10px; font-style: italic; font-size: 0.9em; }}
        .takeaways {{ background: #fff9e6; padding: 15px; margin: 20px 0; border-left: 4px solid #ffd700; }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{course.course_title}</h1>
        <h2>Модуль {module.module_number}: {module.module_title}</h2>
        <h3>{lesson.lesson_title}</h3>
        <p><strong>Цель:</strong> {lesson.lesson_goal}</p>
        <p><strong>Формат:</strong> {lesson.format} | <strong>Время:</strong> {lesson.estimated_time_minutes} мин</p>
    </div>
"""
        
        if content_data.get('learning_objectives'):
            html += """    <div class="objectives">
        <strong>Цели обучения:</strong>
        <ul>
"""
            for obj in content_data['learning_objectives']:
                html += f"            <li>{obj}</li>\n"
            html += """        </ul>
    </div>
"""
        
        slides = content_data.get('slides', [])
        for j, slide in enumerate(slides, 1):
            slide_title = slide.get('slide_title') or slide.get('title', 'Без названия')
            slide_content = slide.get('slide_content') or slide.get('content', '')
            
            html += f"""    <div class="slide">
        <h3>Слайд {j}: {slide_title}</h3>
"""
            
            if slide_content:
                html += f"        <p>{slide_content.replace(chr(10), '<br>')}</p>\n"
            
            if slide.get('code_example'):
                html += f"""        <div class="code">
            <pre><code>{slide['code_example']}</code></pre>
        </div>
"""
            
            if slide.get('notes'):
                html += f"""        <div class="notes">
            <strong>📝 Заметки:</strong> {slide['notes']}
        </div>
"""
            
            html += "    </div>\n"
        
        if content_data.get('key_takeaways'):
            html += """    <div class="takeaways">
        <strong>Ключевые выводы:</strong>
        <ul>
"""
            for key in content_data['key_takeaways']:
                html += f"            <li>{key}</li>\n"
            html += """        </ul>
    </div>
"""
        
        html += """
</body>
</html>"""
        
        return html
    
    @staticmethod
    def export_lesson_pptx(course: Course, module: Module, lesson, content_data: dict) -> bytes:
        """Генерирует PowerPoint для урока"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Титульный слайд
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(102, 126, 234)
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2.5))
            title_frame = title_box.text_frame
            title_frame.text = f"{course.course_title}\n\n{module.module_title}\n\n{lesson.lesson_title}"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.alignment = PP_ALIGN.CENTER
            
            # Цели обучения
            if content_data.get('learning_objectives'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = "Цели обучения"
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.text = ""
                for obj in content_data['learning_objectives']:
                    p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
                    p.text = obj
                    p.font.size = Pt(20)
            
            # Слайды контента
            slides_data = content_data.get('slides', [])
            for slide_data in slides_data:
                slide_title = slide_data.get('slide_title') or slide_data.get('title', 'Без названия')
                slide_content = slide_data.get('slide_content') or slide_data.get('content', '')
                code_example = slide_data.get('code_example')
                
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = slide_title
                
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.word_wrap = True
                
                if slide_content:
                    parts = slide_content.replace('\\n', '\n').split('\n')
                    for i, part in enumerate(parts):
                        part = part.strip()
                        if not part:
                            continue
                        if i == 0:
                            tf.text = part
                            tf.paragraphs[0].font.size = Pt(18)
                        else:
                            p = tf.add_paragraph()
                            p.text = part[2:] if part.startswith('- ') else part
                            p.font.size = Pt(16)
                            if part.startswith('- '):
                                p.level = 1
                
                if code_example:
                    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2.5))
                    code_frame = code_box.text_frame
                    code_frame.text = code_example.replace('\\n', '\n')
                    code_frame.paragraphs[0].font.name = 'Courier New'
                    code_frame.paragraphs[0].font.size = Pt(14)
                    code_box.fill.solid()
                    code_box.fill.fore_color.rgb = RGBColor(40, 44, 52)
                    code_frame.paragraphs[0].font.color.rgb = RGBColor(171, 178, 191)
            
            # Ключевые выводы
            if content_data.get('key_takeaways'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = "Ключевые выводы"
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.text = ""
                for key in content_data['key_takeaways']:
                    p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
                    p.text = key
                    p.font.size = Pt(20)
                    p.font.bold = True
            
            pptx_bytes_io = BytesIO()
            prs.save(pptx_bytes_io)
            pptx_bytes_io.seek(0)
            return pptx_bytes_io.getvalue()
            
        except Exception as e:
            logger.error(f"Ошибка генерации PPTX урока: {e}")
            raise


# Глобальный экземпляр
export_service = ExportService()

