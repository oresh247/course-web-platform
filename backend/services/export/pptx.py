from io import BytesIO
from backend.models.domain import Course, Module


def export_course_pptx(course: Course) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(102, 126, 234)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = course.course_title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_text = f"Целевая аудитория: {course.target_audience}"
    if course.duration_weeks:
        subtitle_text += f" | {course.duration_weeks} недель"
    subtitle_frame.text = subtitle_text
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Структура курса"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = f"Курс состоит из {len(course.modules)} модулей:"
    for module in course.modules:
        p = tf.add_paragraph()
        p.text = f"Модуль {module.module_number}: {module.module_title}"
        p.level = 1
        p.font.size = Pt(18)

    for module in course.modules:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = f"Модуль {module.module_number}: {module.module_title}"
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = f"Цель: {module.module_goal}"
        p = tf.add_paragraph()
        p.text = f"\nУроков в модуле: {len(module.lessons)}"
        p.font.size = Pt(18)

    pptx_bytes_io = BytesIO()
    prs.save(pptx_bytes_io)
    pptx_bytes_io.seek(0)
    return pptx_bytes_io.getvalue()


def export_module_pptx(course: Course, module: Module, content_data: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(102, 126, 234)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = f"{course.course_title}\n\nМодуль {module.module_number}: {module.module_title}"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    goal_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    goal_frame = goal_box.text_frame
    goal_frame.text = f"Цель: {module.module_goal}"
    goal_para = goal_frame.paragraphs[0]
    goal_para.font.size = Pt(18)
    goal_para.font.color.rgb = RGBColor(255, 255, 255)
    goal_para.alignment = PP_ALIGN.CENTER

    lectures = content_data.get('lectures', [])
    for lecture_idx, lecture in enumerate(lectures, 1):
        lecture_title = lecture.get('lecture_title', f'Лекция {lecture_idx}')
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = f"Лекция {lecture_idx}: {lecture_title}"
        if lecture.get('learning_objectives'):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Цели обучения"
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.text = ""
            for obj in lecture['learning_objectives']:
                p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
                p.text = obj
                p.level = 0
                p.font.size = Pt(20)
                p.space_before = Pt(6)
        slides_data = lecture.get('slides', [])
        for slide_data in slides_data:
            slide_title = slide_data.get('slide_title') or slide_data.get('title', 'Без названия')
            slide_content = slide_data.get('slide_content') or slide_data.get('content', '')
            code_example = slide_data.get('code_example')
            if code_example:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = slide_title
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.word_wrap = True
                if slide_content:
                    tf.text = slide_content.replace('\n', '\n')
                    tf.paragraphs[0].font.size = Pt(16)
                code_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(3))
                code_frame = code_box.text_frame
                code_frame.text = code_example.replace('\n', '\n')
                code_para = code_frame.paragraphs[0]
                code_para.font.name = 'Courier New'
                code_para.font.size = Pt(14)
                code_box.fill.solid()
                code_box.fill.fore_color.rgb = RGBColor(40, 44, 52)
                code_para.font.color.rgb = RGBColor(171, 178, 191)
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                title.text = slide_title
                if slide_content:
                    content_shape = slide.placeholders[1]
                    tf = content_shape.text_frame
                    tf.word_wrap = True
                    content_parts = slide_content.replace('\n', '\n').split('\n')
                    for i, part in enumerate(content_parts):
                        part = part.strip()
                        if not part:
                            continue
                        if i == 0:
                            tf.text = part
                            tf.paragraphs[0].font.size = Pt(18)
                            tf.paragraphs[0].space_after = Pt(12)
                        else:
                            p = tf.add_paragraph()
                            p.text = part
                            p.font.size = Pt(16)
                            p.space_before = Pt(6)
                            p.space_after = Pt(6)
                            if part.startswith('- '):
                                p.text = part[2:]
                                p.level = 1
        if lecture.get('key_takeaways'):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = f"Ключевые выводы: {lecture_title}"
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.text = ""
            for key in lecture['key_takeaways']:
                p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
                p.text = key
                p.level = 0
                p.font.size = Pt(20)
                p.font.bold = True

    pptx_bytes_io = BytesIO()
    prs.save(pptx_bytes_io)
    pptx_bytes_io.seek(0)
    return pptx_bytes_io.getvalue()


def export_lesson_pptx(course: Course, module: Module, lesson, content_data: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(102, 126, 234)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2.5))
    title_frame = title_box.text_frame
    title_frame.text = f"{course.course_title}\n\n{module.module_title}\n\n{lesson.lesson_title}"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    if content_data.get('learning_objectives'):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Цели обучения"
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.text = ""
        for obj in content_data['learning_objectives']:
            p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
            p.text = obj
            p.font.size = Pt(20)

    slides_data = content_data.get('slides', [])
    for slide_data in slides_data:
        slide_title = slide_data.get('slide_title') or slide_data.get('title', 'Без названия')
        slide_content = slide_data.get('slide_content') or slide_data.get('content', '')
        code_example = slide_data.get('code_example')
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = slide_title
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.word_wrap = True
        if slide_content:
            parts = slide_content.replace('\n', '\n').split('\n')
            for i, part in enumerate(parts):
                part = part.strip()
                if not part:
                    continue
                if i == 0:
                    tf.text = part
                    tf.paragraphs[0].font.size = Pt(18)
                else:
                    p = tf.add_paragraph()
                    p.text = part[2:] if part.startswith('- ') else part
                    p.font.size = Pt(16)
                    if part.startswith('- '):
                        p.level = 1
        if code_example:
            code_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2.5))
            code_frame = code_box.text_frame
            code_frame.text = code_example.replace('\n', '\n')
            code_frame.paragraphs[0].font.name = 'Courier New'
            code_frame.paragraphs[0].font.size = Pt(14)
            code_box.fill.solid()
            code_box.fill.fore_color.rgb = RGBColor(40, 44, 52)
            code_frame.paragraphs[0].font.color.rgb = RGBColor(171, 178, 191)

    if content_data.get('key_takeaways'):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Ключевые выводы"
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.text = ""
        for key in content_data['key_takeaways']:
            p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
            p.text = key
            p.font.size = Pt(20)
            p.font.bold = True

    pptx_bytes_io = BytesIO()
    prs.save(pptx_bytes_io)
    pptx_bytes_io.seek(0)
    return pptx_bytes_io.getvalue()


